"""Format-specific document ingestors.

Each ingestor reads a source format and produces normalized documents
ready for chunking and embedding. Auto-detected by file extension.
"""

from __future__ import annotations

import csv
import io
import json
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Protocol

logger = logging.getLogger(__name__)


@dataclass(slots=True)
class IngestDocument:
    text: str
    doc_id: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)


class Ingestor(Protocol):
    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]: ...


class PlainTextIngestor:
    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        text = path.read_text(encoding="utf-8", errors="replace")
        if not text.strip():
            return []
        return [IngestDocument(
            text=text,
            doc_id=str(path),
            metadata={"source_type": "file", "source_id": str(path), "format": "text"},
        )]


class MarkdownIngestor:
    """Splits markdown by headers, preserving section hierarchy."""

    _HEADER_RE = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        text = path.read_text(encoding="utf-8", errors="replace")
        if not text.strip():
            return []

        sections = self._split_by_headers(text)
        if not sections:
            return [IngestDocument(
                text=text, doc_id=str(path),
                metadata={"source_type": "file", "source_id": str(path), "format": "markdown"},
            )]

        docs = []
        for i, (header, body) in enumerate(sections):
            combined = f"{header}\n{body}" if header else body
            if not combined.strip():
                continue
            docs.append(IngestDocument(
                text=combined,
                doc_id=f"{path}:section:{i}",
                metadata={
                    "source_type": "file", "source_id": str(path),
                    "format": "markdown", "section": header.strip("# ").strip(),
                    "section_index": i,
                },
            ))
        return docs

    def _split_by_headers(self, text: str) -> list[tuple[str, str]]:
        matches = list(self._HEADER_RE.finditer(text))
        if not matches:
            return [("", text)]

        sections = []
        if matches[0].start() > 0:
            sections.append(("", text[:matches[0].start()]))

        for i, m in enumerate(matches):
            end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
            sections.append((m.group(0), text[m.end():end]))

        return sections


class CodeIngestor:
    """Treats code files as single documents with language metadata."""

    _LANG_MAP = {
        ".py": "python", ".js": "javascript", ".ts": "typescript",
        ".go": "go", ".rs": "rust", ".java": "java", ".rb": "ruby",
        ".php": "php", ".c": "c", ".cpp": "cpp", ".cs": "csharp",
        ".swift": "swift", ".kt": "kotlin", ".sh": "bash",
    }

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        text = path.read_text(encoding="utf-8", errors="replace")
        if not text.strip():
            return []
        lang = self._LANG_MAP.get(path.suffix.lower(), "unknown")
        return [IngestDocument(
            text=text,
            doc_id=str(path),
            metadata={
                "source_type": "file", "source_id": str(path),
                "format": "code", "language": lang,
            },
        )]


class CSVIngestor:
    """Each row becomes a document, or groups of rows if large."""

    def ingest(
        self, path: Path, *, max_rows: int = 10000, **kwargs: Any,
    ) -> list[IngestDocument]:
        text = path.read_text(encoding="utf-8", errors="replace")
        reader = csv.DictReader(io.StringIO(text))
        docs = []
        for i, row in enumerate(reader):
            if i >= max_rows:
                break
            row_text = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
            if not row_text.strip():
                continue
            docs.append(IngestDocument(
                text=row_text,
                doc_id=f"{path}:row:{i}",
                metadata={
                    "source_type": "file", "source_id": str(path),
                    "format": "csv", "row_index": i,
                },
            ))
        return docs


class JSONIngestor:
    """Flattens JSON objects into searchable text documents."""

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        text = path.read_text(encoding="utf-8", errors="replace")
        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            return [IngestDocument(
                text=text, doc_id=str(path),
                metadata={"source_type": "file", "source_id": str(path), "format": "json"},
            )]

        if isinstance(data, list):
            return self._ingest_array(path, data)
        return [IngestDocument(
            text=json.dumps(data, indent=2, ensure_ascii=False),
            doc_id=str(path),
            metadata={"source_type": "file", "source_id": str(path), "format": "json"},
        )]

    def _ingest_array(self, path: Path, items: list) -> list[IngestDocument]:
        docs = []
        for i, item in enumerate(items[:10000]):
            text = json.dumps(item, indent=2, ensure_ascii=False) if isinstance(item, dict) else str(item)
            docs.append(IngestDocument(
                text=text,
                doc_id=f"{path}:item:{i}",
                metadata={
                    "source_type": "file", "source_id": str(path),
                    "format": "json", "item_index": i,
                },
            ))
        return docs


class JSONLIngestor:
    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        docs = []
        with open(path, encoding="utf-8", errors="replace") as f:
            for i, line in enumerate(f):
                line = line.strip()
                if not line:
                    continue
                if i >= 10000:
                    break
                docs.append(IngestDocument(
                    text=line,
                    doc_id=f"{path}:line:{i}",
                    metadata={
                        "source_type": "file", "source_id": str(path),
                        "format": "jsonl", "line_index": i,
                    },
                ))
        return docs


class HTMLIngestor:
    """Extracts text from HTML, strips tags."""

    _TAG_RE = re.compile(r"<[^>]+>")
    _WS_RE = re.compile(r"\s+")

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        raw = path.read_text(encoding="utf-8", errors="replace")
        text = self._TAG_RE.sub(" ", raw)
        text = self._WS_RE.sub(" ", text).strip()
        if not text:
            return []
        return [IngestDocument(
            text=text,
            doc_id=str(path),
            metadata={"source_type": "file", "source_id": str(path), "format": "html"},
        )]


def _extract_pdf_pages(path: Path) -> list[tuple[int, str]]:
    """Extract per-page text from a PDF using pymupdf.

    Returns a list of (page_number_1based, text) tuples. Empty list on
    failure or when pymupdf is not installed. The pdf module used to
    own this code, but it was the only consumer and the wrapper added
    no value, so it was inlined here when the pdf module was removed.
    """
    try:
        import pymupdf
    except ImportError:
        logger.warning(
            "pymupdf not installed - PDF ingestion disabled. "
            "Install with: pip install pymupdf",
        )
        return []
    try:
        doc = pymupdf.open(str(path))
        pages = [(idx + 1, doc[idx].get_text("text")) for idx in range(len(doc))]
        doc.close()
        return pages
    except Exception as exc:
        logger.warning("PDF read failed for %s: %s", path, exc)
        return []


class PDFIngestor:
    """Read a PDF file and produce one IngestDocument per page."""

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        # Sync entry point - used by the synchronous IndexingEngine path.
        # The async path runs the same extraction in a thread.
        return self._build_docs(path, _extract_pdf_pages(path))

    async def ingest_async(self, path: Path, bus: Any = None) -> list[IngestDocument]:
        # ``bus`` is accepted for API compatibility with sibling ingestors
        # that still delegate over ServiceBus; PDFIngestor reads directly.
        import asyncio
        pages = await asyncio.to_thread(_extract_pdf_pages, path)
        return self._build_docs(path, pages)

    def _build_docs(
        self, path: Path, pages: list[tuple[int, str]],
    ) -> list[IngestDocument]:
        docs: list[IngestDocument] = []
        for page_num, text in pages:
            text = text.strip()
            if not text:
                continue
            docs.append(IngestDocument(
                text=text,
                doc_id=f"{path}:page:{page_num}",
                metadata={
                    "source_type": "file",
                    "source_id": str(path),
                    "format": "pdf",
                    "page": page_num,
                },
            ))
        return docs


def _extract_docx_sections(path: Path) -> list[tuple[str, str]]:
    """Extract section-aligned text from a .docx using python-docx.

    Returns a list of ``(heading, body)`` tuples. ``heading`` is the
    heading paragraph's text (empty for the content above the first
    heading); ``body`` is the joined non-empty paragraphs that follow
    until the next heading. Tables are flattened to ``cell | cell``
    rows so structured content still reaches the index.

    Empty list on failure or when python-docx is not installed.
    Mirrors :func:`_extract_pdf_pages` and
    :func:`_extract_spreadsheet_sheets` so the calling ingestor stays
    a thin wrapper.
    """
    try:
        import docx  # python-docx
    except ImportError:
        logger.warning(
            "python-docx not installed - DOCX ingestion disabled. "
            "Install with: pip install python-docx",
        )
        return []
    try:
        doc = docx.Document(str(path))
    except Exception as exc:
        logger.warning("DOCX read failed for %s: %s", path, exc)
        return []

    sections: list[tuple[str, list[str]]] = []
    current_heading = ""
    current_body: list[str] = []

    def _flush() -> None:
        if current_heading or current_body:
            sections.append((current_heading, list(current_body)))

    # Walk paragraphs in document order. Heading-styled paragraphs
    # split sections; everything else accumulates. python-docx exposes
    # tables separately, so we sweep them after - keeps the helper
    # tolerant of malformed docs that mix tables and paragraphs
    # heavily.
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style_name = ""
        try:
            style_name = (para.style.name if para.style else "") or ""
        except Exception:
            pass
        if style_name.startswith("Heading"):
            _flush()
            current_heading = text
            current_body = []
        else:
            current_body.append(text)
    _flush()

    # Flatten any table cells into the same bucket as the last
    # section, or a dedicated "Tables" section if there's nothing
    # before them.
    for table in getattr(doc, "tables", []):
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            row_text = " | ".join(c for c in cells if c)
            if not row_text:
                continue
            if not sections:
                sections.append(("Tables", []))
            sections[-1][1].append(row_text)

    return [
        (heading, "\n".join(body))
        for heading, body in sections
        if heading or body
    ]


def _extract_spreadsheet_sheets(
    path: Path, max_rows: int = 10000,
) -> list[tuple[str, list[Any], list[list[Any]]]]:
    """Extract (sheet_name, headers, rows) tuples from a spreadsheet.

    Supports .csv (stdlib csv) and .xlsx (openpyxl, optional). Returns
    an empty list on unsupported formats or missing libraries. The
    spreadsheet module used to own this code, but it was the only
    consumer and was inlined when the module was removed.
    """
    suffix = path.suffix.lower()
    if suffix == ".csv":
        import csv
        try:
            with path.open(encoding="utf-8", newline="") as f:
                reader = csv.reader(f)
                rows = list(reader)
            if not rows:
                return []
            headers = rows[0]
            return [("Sheet1", headers, rows[1:max_rows + 1])]
        except Exception as exc:
            logger.warning("CSV read failed for %s: %s", path, exc)
            return []
    if suffix in (".xlsx", ".xlsm"):
        try:
            import openpyxl
        except ImportError:
            logger.warning(
                "openpyxl not installed - XLSX ingestion disabled. "
                "Install with: pip install openpyxl",
            )
            return []
        try:
            wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
            out: list[tuple[str, list[Any], list[list[Any]]]] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows_iter = ws.iter_rows(values_only=True)
                first = next(rows_iter, None)
                if first is None:
                    continue
                headers = [c if c is not None else "" for c in first]
                rows: list[list[Any]] = []
                for i, row in enumerate(rows_iter):
                    if i >= max_rows:
                        break
                    rows.append(list(row))
                out.append((sheet_name, headers, rows))
            wb.close()
            return out
        except Exception as exc:
            logger.warning("XLSX read failed for %s: %s", path, exc)
            return []
    return []


def _extract_pptx_slides(path: Path) -> list[tuple[int, str]]:
    """Extract per-slide text from a .pptx using python-pptx.

    Returns ``[(slide_number_1based, text), ...]``. Empty list on
    failure or when python-pptx is missing. Mirrors the
    ``_extract_pdf_pages`` contract exactly so the wrapping ingestor
    stays trivial.
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning(
            "python-pptx not installed - PPTX ingestion disabled. "
            "Install with: pip install python-pptx",
        )
        return []
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        logger.warning("PPTX read failed for %s: %s", path, exc)
        return []

    out: list[tuple[int, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = getattr(shape, "text_frame", None)
            if tf is None:
                continue
            text = (tf.text or "").strip()
            if text:
                chunks.append(text)
        # Slide notes routinely carry the speaker's talking points -
        # often more semantic content than the bullets on screen.
        try:
            notes = getattr(slide, "notes_slide", None)
            if notes is not None:
                ntf = getattr(notes, "notes_text_frame", None)
                if ntf is not None:
                    note_text = (ntf.text or "").strip()
                    if note_text:
                        chunks.append(f"[notes] {note_text}")
        except Exception:
            pass
        joined = "\n".join(chunks).strip()
        if joined:
            out.append((idx, joined))
    return out


def _extract_odt_paragraphs(path: Path) -> list[tuple[str, str]]:
    """Extract section-aligned text from a .odt by reading its
    ``content.xml`` directly (no odfpy dep).

    Returns ``[(heading, body), ...]`` matching ``_extract_docx_sections``
    so a future generic "sectioned" ingestor base can subsume both.
    Empty list on read / parse failure - the rag fallback path then
    treats the file as opaque text.
    """
    import zipfile
    from xml.etree import ElementTree as ET

    try:
        with zipfile.ZipFile(str(path)) as zf:
            with zf.open("content.xml") as fh:
                raw = fh.read().decode("utf-8", errors="replace")
    except Exception as exc:
        logger.warning("ODT read failed for %s: %s", path, exc)
        return []

    try:
        root = ET.fromstring(raw)
    except ET.ParseError as exc:
        logger.warning("ODT XML parse failed for %s: %s", path, exc)
        return []

    # ODF namespaces. We never bind them by prefix because XML wants
    # full ``{uri}local`` syntax under ElementTree.
    NS_TEXT = "{urn:oasis:names:tc:opendocument:xmlns:text:1.0}"
    sections: list[tuple[str, list[str]]] = []
    current_heading = ""
    current_body: list[str] = []

    def _flush() -> None:
        if current_heading or current_body:
            sections.append((current_heading, list(current_body)))

    def _text_of(node: Any) -> str:
        # ODF wraps inline runs in <text:span> - join all descendant
        # text the cheap way via itertext().
        return "".join(node.itertext()).strip()

    for node in root.iter():
        tag = node.tag
        if tag == f"{NS_TEXT}h":
            _flush()
            current_heading = _text_of(node)
            current_body = []
        elif tag == f"{NS_TEXT}p":
            txt = _text_of(node)
            if txt:
                current_body.append(txt)
    _flush()

    return [
        (heading, "\n".join(body))
        for heading, body in sections
        if heading or body
    ]


def _extract_ods_sheets(
    path: Path, max_rows: int = 10000,
) -> list[tuple[str, list[Any], list[list[Any]]]]:
    """Extract (sheet_name, headers, rows) from an .ods via raw XML.

    Same return shape as :func:`_extract_spreadsheet_sheets` so the
    existing :class:`SpreadsheetIngestor` can swallow it without
    changes - we just register a new sync entry that points at this
    helper indirectly through :class:`ODSIngestor` below.
    """
    import zipfile
    from xml.etree import ElementTree as ET

    try:
        with zipfile.ZipFile(str(path)) as zf:
            with zf.open("content.xml") as fh:
                raw = fh.read().decode("utf-8", errors="replace")
    except Exception as exc:
        logger.warning("ODS read failed for %s: %s", path, exc)
        return []

    try:
        root = ET.fromstring(raw)
    except ET.ParseError as exc:
        logger.warning("ODS XML parse failed for %s: %s", path, exc)
        return []

    NS_TABLE = "{urn:oasis:names:tc:opendocument:xmlns:table:1.0}"
    out: list[tuple[str, list[Any], list[list[Any]]]] = []

    for tbl in root.iter(f"{NS_TABLE}table"):
        name = tbl.attrib.get(f"{NS_TABLE}name") or "Sheet"
        rows: list[list[Any]] = []
        for row in tbl.iter(f"{NS_TABLE}table-row"):
            if len(rows) >= max_rows + 1:
                break
            cells: list[str] = []
            for cell in row.iter(f"{NS_TABLE}table-cell"):
                txt = "".join(cell.itertext()).strip()
                cells.append(txt)
            # Trim trailing empties (ODS pads rows with empty cells).
            while cells and not cells[-1]:
                cells.pop()
            if cells:
                rows.append(cells)
        if not rows:
            continue
        headers = rows[0]
        out.append((name, headers, rows[1:max_rows + 1]))
    return out


def _extract_rtf_text(path: Path) -> str:
    """Strip a .rtf to plain text without a heavy dependency.

    RTF is a small grammar: control words ``\\foo``, hex escapes
    ``\\'XX``, group braces ``{}``, deletable groups ``{\\*...}``.
    A regex pass covers the long tail of business documents well
    enough for retrieval (we lose formatting but keep words and
    paragraph breaks). Returns ``""`` on read failure.
    """
    try:
        raw = path.read_text(encoding="latin-1", errors="replace")
    except Exception as exc:
        logger.warning("RTF read failed for %s: %s", path, exc)
        return ""

    # Drop ``{\*\anything ...}`` blocks - they hold metadata the LLM
    # doesn't need (fonts, colours, revision history).
    def _strip_deletable(text: str) -> str:
        out: list[str] = []
        depth = 0
        i = 0
        n = len(text)
        while i < n:
            if text[i] == "{" and i + 1 < n and text[i + 1] == "\\" and i + 2 < n and text[i + 2] == "*":
                depth += 1
                i += 1
                continue
            if depth > 0:
                if text[i] == "{":
                    depth += 1
                elif text[i] == "}":
                    depth -= 1
                i += 1
                continue
            out.append(text[i])
            i += 1
        return "".join(out)

    text = _strip_deletable(raw)
    # Hex escapes ``\'XX`` → byte. Decode as cp1252 (Windows default).
    def _hex_sub(m: re.Match[str]) -> str:
        try:
            return bytes([int(m.group(1), 16)]).decode("cp1252", errors="replace")
        except Exception:
            return ""
    text = re.sub(r"\\'([0-9a-fA-F]{2})", _hex_sub, text)
    # Common control words → newline / space / nothing.
    text = re.sub(r"\\par\b ?", "\n", text)
    text = re.sub(r"\\line\b ?", "\n", text)
    text = re.sub(r"\\tab\b ?", "\t", text)
    # Generic control words: ``\foo123 `` or ``\foo`` followed by space.
    text = re.sub(r"\\[a-zA-Z]+-?\d*\s?", "", text)
    # Drop stray braces.
    text = text.replace("{", "").replace("}", "")
    # Collapse whitespace runs but preserve paragraph breaks.
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n[ \t]+", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


class PPTXIngestor:
    """One IngestDocument per slide. Mirrors PDFIngestor verbatim so
    the rag pipeline can't tell PPTX apart from PDF at the metadata
    layer - both produce per-page citations via the existing
    ``[file · page N]`` rendering.
    """

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        slides = _extract_pptx_slides(path)
        if not slides:
            return []
        docs: list[IngestDocument] = []
        for slide_num, text in slides:
            t = (text or "").strip()
            if not t:
                continue
            docs.append(IngestDocument(
                text=t,
                doc_id=f"{path}:slide:{slide_num}",
                metadata={
                    "source_type": "file",
                    "source_id": str(path),
                    "format": "pptx",
                    # ``page`` so the existing citation renderer
                    # (``_format_rag_context_block``) emits
                    # ``[deck.pptx · page 4]`` without any change.
                    "page": slide_num,
                },
            ))
        return docs


class ODTIngestor:
    """ODT mirror of DOCXIngestor. Section-aligned, header-driven."""

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        sections = _extract_odt_paragraphs(path)
        if not sections:
            return []
        docs: list[IngestDocument] = []
        for i, (heading, body) in enumerate(sections):
            text_parts = [p for p in (heading, body) if p]
            text = "\n".join(text_parts).strip()
            if not text:
                continue
            docs.append(IngestDocument(
                text=text,
                doc_id=f"{path}:section:{i}",
                metadata={
                    "source_type": "file",
                    "source_id": str(path),
                    "format": "odt",
                    "section": heading,
                    "section_index": i,
                },
            ))
        return docs


class ODSIngestor:
    """ODS mirror of SpreadsheetIngestor - one IngestDocument per row,
    sheet name in metadata. Delegates extraction to
    :func:`_extract_ods_sheets` which returns the same shape as the
    XLSX extractor so the row-formatting loop is identical."""

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        sheets = _extract_ods_sheets(path)
        if not sheets:
            return []
        docs: list[IngestDocument] = []
        for sheet_name, headers, rows in sheets:
            for i, row in enumerate(rows):
                if headers:
                    text = " | ".join(
                        f"{h}: {row[j]}"
                        for j, h in enumerate(headers)
                        if j < len(row) and row[j]
                    )
                else:
                    text = " | ".join(str(v) for v in row if v)
                if not text.strip():
                    continue
                docs.append(IngestDocument(
                    text=text,
                    doc_id=f"{path}:{sheet_name}:row:{i}",
                    metadata={
                        "source_type": "file",
                        "source_id": str(path),
                        "format": "ods",
                        "sheet": sheet_name,
                        "row_index": i,
                    },
                ))
        return docs


class RTFIngestor:
    """RTF → plain text single-doc ingestor. The text loses formatting
    but keeps the words, which is all retrieval needs."""

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        text = _extract_rtf_text(path)
        if not text:
            return []
        return [IngestDocument(
            text=text,
            doc_id=str(path),
            metadata={
                "source_type": "file",
                "source_id": str(path),
                "format": "rtf",
            },
        )]


class DOCXIngestor:
    """Read a .docx (Word document) and produce one IngestDocument
    per section (heading + body). Mirrors PDFIngestor's shape so the
    rag pipeline's per-doc metadata path stays uniform.

    No async variant: python-docx is fast enough that the sync
    ingestor path off-loaded via ``asyncio.to_thread`` in the rag
    module is sufficient. If we ever hit a 500-page DOCX that needs
    streaming, mirror PDFIngestor.ingest_async then.
    """

    def ingest(self, path: Path, **kwargs: Any) -> list[IngestDocument]:
        sections = _extract_docx_sections(path)
        if not sections:
            return []
        docs: list[IngestDocument] = []
        for i, (heading, body) in enumerate(sections):
            text_parts = [p for p in (heading, body) if p]
            text = "\n".join(text_parts).strip()
            if not text:
                continue
            docs.append(IngestDocument(
                text=text,
                doc_id=f"{path}:section:{i}",
                metadata={
                    "source_type": "file",
                    "source_id": str(path),
                    "format": "docx",
                    "section": heading,
                    "section_index": i,
                },
            ))
        return docs


class SpreadsheetIngestor:
    """Read a spreadsheet (.csv / .xlsx) and produce one IngestDocument per row."""

    async def ingest_async(self, path: Path, bus: Any = None) -> list[IngestDocument]:
        import asyncio
        sheets = await asyncio.to_thread(_extract_spreadsheet_sheets, path)
        docs: list[IngestDocument] = []
        for sheet_name, headers, rows in sheets:
            for i, row in enumerate(rows):
                if headers:
                    text = " | ".join(
                        f"{h}: {row[j]}"
                        for j, h in enumerate(headers)
                        if j < len(row) and row[j]
                    )
                else:
                    text = " | ".join(str(v) for v in row if v)
                if not text.strip():
                    continue
                docs.append(IngestDocument(
                    text=text,
                    doc_id=f"{path}:{sheet_name}:row:{i}",
                    metadata={
                        "source_type": "file",
                        "source_id": str(path),
                        "format": "spreadsheet",
                        "sheet": sheet_name,
                        "row_index": i,
                    },
                ))
        return docs


# ── Registry ──────────────────────────────────────────────────────────

_SYNC_INGESTORS: dict[str, Ingestor] = {
    ".md": MarkdownIngestor(),
    ".markdown": MarkdownIngestor(),
    ".txt": PlainTextIngestor(),
    ".rst": PlainTextIngestor(),
    ".log": PlainTextIngestor(),
    ".py": CodeIngestor(),
    ".js": CodeIngestor(),
    ".ts": CodeIngestor(),
    ".tsx": CodeIngestor(),
    ".jsx": CodeIngestor(),
    ".go": CodeIngestor(),
    ".rs": CodeIngestor(),
    ".java": CodeIngestor(),
    ".rb": CodeIngestor(),
    ".php": CodeIngestor(),
    ".c": CodeIngestor(),
    ".cpp": CodeIngestor(),
    ".h": CodeIngestor(),
    ".cs": CodeIngestor(),
    ".swift": CodeIngestor(),
    ".kt": CodeIngestor(),
    ".sh": CodeIngestor(),
    ".bash": CodeIngestor(),
    ".yaml": PlainTextIngestor(),
    ".yml": PlainTextIngestor(),
    ".toml": PlainTextIngestor(),
    ".ini": PlainTextIngestor(),
    ".cfg": PlainTextIngestor(),
    ".conf": PlainTextIngestor(),
    ".csv": CSVIngestor(),
    ".tsv": CSVIngestor(),
    ".json": JSONIngestor(),
    ".jsonl": JSONLIngestor(),
    ".ndjson": JSONLIngestor(),
    ".html": HTMLIngestor(),
    ".htm": HTMLIngestor(),
    ".xml": PlainTextIngestor(),
    ".sql": CodeIngestor(),
    ".docx": DOCXIngestor(),
    ".pptx": PPTXIngestor(),
    ".odt": ODTIngestor(),
    ".ods": ODSIngestor(),
    ".rtf": RTFIngestor(),
}

_ASYNC_EXTENSIONS = {".pdf", ".xlsx", ".xls"}


def get_ingestor(ext: str) -> Ingestor | None:
    return _SYNC_INGESTORS.get(ext.lower())


def is_async_format(ext: str) -> bool:
    return ext.lower() in _ASYNC_EXTENSIONS


def supported_extensions() -> set[str]:
    return set(_SYNC_INGESTORS.keys()) | _ASYNC_EXTENSIONS
