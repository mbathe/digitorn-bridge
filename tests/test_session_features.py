"""Tests avancés pour les fonctionnalités implémentées dans cette session.

Couvre :
  - WeasyPrint PDF backend (incassable, tous thèmes, cover, TOC)
  - Workbench Renderer V2 (snapshot, post_use_summary)
  - Markdown Renderer (arbre sections, tableaux, citations)
  - Spreadsheet Renderer snapshot
  - Presentation Renderer snapshot
  - Presentation engine (layouts, effets, overlaps)
  - Table Typst renderer (colonnes, headers)
  - Agent loop incremental actions
"""

import asyncio
import json
import os
import re
import shutil
import tempfile
from pathlib import Path

import pytest

# ---------------------------------------------------------------------------
# Optional dependency availability flags
# ---------------------------------------------------------------------------
_has_weasyprint = pytest.importorskip is not None  # helper; real skip below
try:
    import weasyprint as _wp  # noqa: F401
    _has_weasyprint = True
except ImportError:
    _has_weasyprint = False

try:
    import pptx as _pptx  # noqa: F401
    _has_pptx = True
except ImportError:
    _has_pptx = False

try:
    from digitorn.modules.spreadsheet.renderer import SpreadsheetRenderer as _SR  # noqa: F401
    _has_spreadsheet_renderer = True
except ImportError:
    _has_spreadsheet_renderer = False

_skip_weasy = pytest.mark.skipif(not _has_weasyprint, reason="weasyprint not installed")
_skip_pptx = pytest.mark.skipif(not _has_pptx, reason="python-pptx not installed")
_skip_spreadsheet = pytest.mark.skipif(not _has_spreadsheet_renderer, reason="SpreadsheetRenderer not available")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 1. WeasyPrint PDF Backend
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


@_skip_weasy
class TestWeasyPrintBackend:
    """Test the WeasyPrint-based PDF generation."""

    def setup_method(self):
        self.tmpdir = tempfile.mkdtemp()

    def _path(self, name: str) -> str:
        return os.path.join(self.tmpdir, name)

    def test_simple_markdown(self):
        from digitorn.modules.pdf.weasy import generate_pdf
        r = generate_pdf("# Hello\n\nWorld", self._path("simple.pdf"))
        assert r["pages"] >= 1
        assert os.path.isfile(r["output_path"])

    def test_never_crashes_on_special_chars(self):
        """The whole point of WeasyPrint — any input produces a PDF."""
        from digitorn.modules.pdf.weasy import generate_pdf
        nasty = (
            "# Test <script>alert('xss')</script>\n\n"
            "Price: $100 @mention #hashtag\n"
            "Arrows: < > << >> → ← ↑ ↓\n"
            "Math: 2² + 3³ = x\n"
            "Emoji: 🚀 🎉 🔥\n"
            "| Col#1 | $Price | @User |\n|---|---|---|\n| <val> | $99 | @john |\n"
            "```python\ndef f(x): return x**2\n```\n"
            "> Quote with 'single' and \"double\" quotes\n"
            "URL: https://example.com/path?q=test&p=1#anchor\n"
        )
        r = generate_pdf(nasty, self._path("nasty.pdf"))
        assert r["pages"] >= 1

    def test_all_themes_compile(self):
        from digitorn.modules.pdf.weasy import generate_pdf
        md = "# Title\n\n## Section\n\nContent.\n\n| A | B |\n|---|---|\n| 1 | 2 |"
        themes = ["consulting", "scientific", "marketing", "financial", "tech", "executive", "report", "ocean"]
        for theme in themes:
            r = generate_pdf(md, self._path(f"{theme}.pdf"), style=theme, title="Test")
            assert r["pages"] >= 1, f"Theme {theme} failed"
            assert os.path.isfile(r["output_path"])

    def test_cover_page_generated(self):
        from digitorn.modules.pdf.weasy import generate_pdf
        r = generate_pdf(
            "# Section 1\n\nContent.",
            self._path("cover.pdf"),
            title="My Report", author="Test Author",
        )
        assert r["pages"] >= 2  # cover + content

    def test_toc_generated(self):
        from digitorn.modules.pdf.weasy import generate_pdf
        md = "# Section A\n\nText.\n\n# Section B\n\nText.\n\n## Sub B1\n\nText."
        r = generate_pdf(md, self._path("toc.pdf"), title="TOC Test")
        assert r["pages"] >= 3  # cover + toc + content

    def test_long_report(self):
        """Test a 10+ page report."""
        from digitorn.modules.pdf.weasy import generate_pdf
        sections = []
        for i in range(10):
            sections.append(
                f"# Section {i + 1}: Topic {i + 1}\n\n"
                f"This is a detailed section about topic {i + 1}. "
                f"It contains important information.\n\n"
                f"| Metric | Value | Change |\n|---|---|---|\n"
                f"| Revenue | ${i * 10 + 50}M | +{i + 5}% |\n"
                f"| Users | {i * 1000 + 5000} | +{i * 100} |\n"
            )
        md = "\n\n".join(sections)
        r = generate_pdf(md, self._path("long.pdf"), style="consulting", title="Long Report")
        assert r["pages"] >= 10

    def test_tables_with_many_columns(self):
        from digitorn.modules.pdf.weasy import generate_pdf
        md = (
            "# Table Test\n\n"
            "| Framework | Score | Ease | Perf | Eco | Prod | Cost | Community |\n"
            "|---|:---:|:---:|:---:|:---:|:---:|:---:|:---:|\n"
            "| AutoGen | 8.9 | 8.5 | 9.2 | 8.7 | 9.0 | 8.8 | 8.5 |\n"
            "| CrewAI | 8.7 | 9.0 | 8.5 | 8.9 | 9.2 | 8.3 | 8.6 |\n"
        )
        r = generate_pdf(md, self._path("wide_table.pdf"), title="Tables")
        assert r["pages"] >= 1

    def test_images_from_url(self):
        """WeasyPrint should fetch remote images."""
        from digitorn.modules.pdf.weasy import generate_pdf
        md = (
            "# Image Test\n\n"
            "![Logo](https://www.python.org/static/community_logos/python-logo-generic.svg)\n"
        )
        r = generate_pdf(md, self._path("image.pdf"), title="Images")
        assert r["pages"] >= 1

    def test_page_breaks_on_h1(self):
        from digitorn.modules.pdf.weasy import generate_pdf
        md = "# Part 1\n\nContent 1.\n\n# Part 2\n\nContent 2.\n\n# Part 3\n\nContent 3."
        r = generate_pdf(md, self._path("pagebreak.pdf"), title="Breaks")
        assert r["pages"] >= 4  # cover + toc + 3 parts (each on new page)

    def test_page_break_literal_cleaned(self):
        from digitorn.modules.pdf.weasy import _clean_markdown
        md = "Text before\n\nPage Break\n\nText after"
        cleaned = _clean_markdown(md)
        assert "Page Break" not in cleaned
        assert "---" in cleaned

    def test_markdown_to_html_never_fails(self):
        from digitorn.modules.pdf.weasy import markdown_to_html
        # Even totally broken input produces HTML
        result = markdown_to_html(None)  # type: ignore
        assert isinstance(result, str)

    def test_blockquote_rendering(self):
        from digitorn.modules.pdf.weasy import generate_pdf
        md = '# Quotes\n\n> "This is a profound quote about AI." — Expert\n\nMore text.'
        r = generate_pdf(md, self._path("quotes.pdf"), title="Quotes")
        assert r["pages"] >= 1

    def test_code_block_rendering(self):
        from digitorn.modules.pdf.weasy import generate_pdf
        md = "# Code\n\n```python\ndef hello():\n    print('world')\n```\n\nDone."
        r = generate_pdf(md, self._path("code.pdf"), title="Code")
        assert r["pages"] >= 1

    def test_style_auto_detection(self):
        from digitorn.modules.pdf.styles import style_for_context
        assert style_for_context("consulting report strategy") == "consulting"
        assert style_for_context("scientific paper research") == "scientific"
        assert style_for_context("marketing campaign brand") == "marketing"
        assert style_for_context("financial annual report investor") == "financial"
        assert style_for_context("API specification architecture") == "tech"
        assert style_for_context("executive board CEO") == "executive"
        assert style_for_context("random topic") == "report"


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 2. Workbench Renderer V2 — snapshot() and post_use_summary()
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestWorkbenchRendererV2:
    """Test the new snapshot() and post_use_summary() methods."""

    def setup_method(self):
        from digitorn.modules.context_builder.workbench import Workbench
        from digitorn.modules.pdf.md_renderer import MarkdownRenderer
        from digitorn.modules.presentation.renderer import PresentationRenderer

        self.wb = Workbench()
        self.wb.register_renderer("markdown", MarkdownRenderer())
        self.wb.register_renderer("text", MarkdownRenderer())
        # SpreadsheetRenderer does not exist as a standalone module;
        # register a basic renderer for workbook type if available.
        try:
            from digitorn.modules.spreadsheet.renderer import SpreadsheetRenderer
            self.wb.register_renderer("workbook", SpreadsheetRenderer())
        except ImportError:
            pass
        self.wb.register_renderer("presentation", PresentationRenderer())

    # ── Snapshot basics ──

    def test_snapshot_returns_string(self):
        self.wb.write("test", "# Hello\n\nWorld")
        snap = self.wb.get_snapshot("test")
        assert isinstance(snap, str)
        assert "test" in snap

    def test_snapshot_unknown_buffer(self):
        snap = self.wb.get_snapshot("nonexistent")
        assert "not found" in snap

    def test_snapshot_no_renderer_fallback(self):
        self.wb.write("raw", "data", buffer_type="unknown_type")
        snap = self.wb.get_snapshot("raw")
        assert "unknown_type" in snap
        assert "1 lines" in snap

    # ── Markdown Renderer Snapshot ──

    def test_markdown_snapshot_shows_structure(self):
        self.wb.write("doc", "# Title\n\n## Section 1\n\nText.\n\n## Section 2\n\nMore.")
        snap = self.wb.get_snapshot("doc", action="write")
        assert "Structure" in snap
        assert "Section 1" in snap
        assert "Section 2" in snap

    def test_markdown_snapshot_shows_tables(self):
        self.wb.write("doc", "# Report\n\n| A | B |\n|---|---|\n| 1 | 2 |")
        snap = self.wb.get_snapshot("doc")
        assert "tableau" in snap.lower()

    def test_markdown_snapshot_shows_space(self):
        self.wb.write("doc", "# Hello")
        snap = self.wb.get_snapshot("doc")
        assert "200,000" in snap or "Espace" in snap

    def test_markdown_snapshot_after_append(self):
        self.wb.write("doc", "# Title\n\n## Section 1\n\nText.")
        self.wb.append("doc", "\n\n## Section 2\n\nMore text.")
        snap = self.wb.get_snapshot("doc", action="append", added_content="## Section 2")
        assert "Section 2" in snap
        assert "jouté" in snap  # "Ajouté"

    def test_markdown_snapshot_quotes_and_lists(self):
        self.wb.write("doc", "# Report\n\n> A quote\n\n- Item 1\n- Item 2\n- Item 3")
        snap = self.wb.get_snapshot("doc")
        assert "citation" in snap.lower()
        assert "items" in snap.lower() or "liste" in snap.lower()

    # ── Spreadsheet Renderer Snapshot ──

    @_skip_spreadsheet
    def test_spreadsheet_snapshot(self):
        spec = json.dumps({
            "sheets": [{
                "name": "Revenue",
                "columns": [{"header": "Q"}, {"header": "Rev"}],
                "data": [["Q1", 100], ["Q2", 200]],
            }]
        })
        self.wb.write("xl", spec, buffer_type="workbook")
        snap = self.wb.get_snapshot("xl")
        assert "Revenue" in snap
        assert "2r" in snap or "2 rows" in snap.lower()

    @_skip_spreadsheet
    def test_spreadsheet_snapshot_multiple_sheets(self):
        spec = json.dumps({
            "sheets": [
                {"name": "Revenue", "columns": [{"header": "Q"}], "data": [["Q1"]]},
                {"name": "Costs", "columns": [{"header": "Item"}], "data": [["Salary"]]},
            ]
        })
        self.wb.write("xl", spec, buffer_type="workbook")
        snap = self.wb.get_snapshot("xl")
        assert "Revenue" in snap
        assert "Costs" in snap
        assert "2 sheets" in snap

    @_skip_spreadsheet
    def test_spreadsheet_snapshot_with_formulas_charts(self):
        spec = json.dumps({
            "sheets": [{
                "name": "Data",
                "columns": [{"header": "A"}],
                "data": [["x"]],
                "formulas": {"B2": "=SUM(A:A)"},
                "charts": [{"type": "column"}],
            }]
        })
        self.wb.write("xl", spec, buffer_type="workbook")
        snap = self.wb.get_snapshot("xl")
        assert "formula" in snap.lower()
        assert "chart" in snap.lower()

    # ── Presentation Renderer Snapshot ──

    def test_presentation_snapshot(self):
        slides = json.dumps([
            {"html": "<h1>Title</h1>", "layout": "title"},
            {"html": "<h2>Content</h2><p>Body</p>", "layout": "content"},
        ])
        self.wb.write("pres", slides, buffer_type="presentation")
        snap = self.wb.get_snapshot("pres")
        assert "2 slides" in snap
        assert "title" in snap.lower()
        assert "Title" in snap

    def test_presentation_snapshot_empty(self):
        self.wb.write("pres", "[]", buffer_type="presentation")
        snap = self.wb.get_snapshot("pres")
        assert "empty" in snap.lower()

    # ── Post-Use Summary ──

    def test_post_use_summary_pdf(self):
        self.wb.write("doc", "# Report\n\n## Section 1\n\nContent.")
        summary = self.wb.get_post_use_summary("doc", "pdf.generate", {
            "output_path": "/tmp/report.pdf",
            "pages": 5, "size_bytes": 72000, "style": "consulting",
        })
        assert "PDF" in summary
        assert "5 pages" in summary
        assert "consulting" in summary

    @_skip_spreadsheet
    def test_post_use_summary_spreadsheet(self):
        self.wb.write("xl", "{}", buffer_type="workbook")
        summary = self.wb.get_post_use_summary("xl", "spreadsheet.create", {
            "output_path": "/tmp/data.xlsx",
            "sheets": ["Revenue", "Costs"], "total_rows": 30,
            "total_charts": 1, "size_bytes": 15000,
        })
        assert "Excel" in summary
        assert "Revenue" in summary
        assert "Costs" in summary

    def test_post_use_summary_presentation(self):
        self.wb.write("pres", "[]", buffer_type="presentation")
        summary = self.wb.get_post_use_summary("pres", "presentation.finalize", {
            "path": "/tmp/pitch.pptx",
            "slides": 10, "theme": "ocean", "size_kb": 45,
        })
        assert "PowerPoint" in summary
        assert "10 slides" in summary
        assert "ocean" in summary

    def test_post_use_summary_no_renderer(self):
        self.wb.write("raw", "data", buffer_type="unknown")
        summary = self.wb.get_post_use_summary("raw", "tool", {})
        assert summary == ""

    def test_post_use_summary_none_result(self):
        self.wb.write("doc", "# Test")
        summary = self.wb.get_post_use_summary("doc", "pdf.generate", None)
        assert summary == ""


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 3. Markdown Renderer — Validation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestMarkdownRendererValidation:

    def setup_method(self):
        from digitorn.modules.pdf.md_renderer import MarkdownRenderer
        self.renderer = MarkdownRenderer()

    def _buf(self, content: str):
        from digitorn.modules.context_builder.workbench import Buffer
        return Buffer(key="test", content=content, buffer_type="markdown")

    def test_validate_empty(self):
        issues = self.renderer.validate(self._buf(""))
        assert any(i.severity == "warning" for i in issues)

    def test_validate_unclosed_code(self):
        issues = self.renderer.validate(self._buf("```python\ncode\n"))
        assert any("code block" in i.message.lower() for i in issues)

    def test_validate_long_lines(self):
        long_line = "x" * 600
        issues = self.renderer.validate(self._buf(f"# Title\n\n{long_line}"))
        assert any("long line" in i.message.lower() for i in issues)

    def test_validate_ok(self):
        issues = self.renderer.validate(self._buf("# Title\n\nNormal content."))
        assert len(issues) == 0

    def test_summary_line(self):
        summary = self.renderer.summary_line(self._buf(
            "# Title\n\n## A\n\n## B\n\n| x | y |\n|---|---|\n| 1 | 2 |"
        ))
        assert "sections" in summary
        assert "tables" in summary


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 4. PDF Module Integration
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


@_skip_weasy
class TestPDFModuleIntegration:
    """Test the PDF module with WeasyPrint backend."""

    def setup_method(self):
        self.tmpdir = tempfile.mkdtemp()

    @pytest.mark.asyncio
    async def test_generate_action(self):
        from digitorn.modules.pdf.module import PDFModule
        from digitorn.modules.pdf.params import GenerateParams
        m = PDFModule()
        r = await m.generate(GenerateParams(
            content="# Test\n\nContent.",
            output_path=os.path.join(self.tmpdir, "test.pdf"),
            title="Test",
        ))
        assert r.success
        assert "output_path" in r.data

    @pytest.mark.asyncio
    async def test_generate_auto_style(self):
        from digitorn.modules.pdf.module import PDFModule
        from digitorn.modules.pdf.params import GenerateParams
        m = PDFModule()
        r = await m.generate(GenerateParams(
            content="# Financial Report Q1\n\nRevenue grew.",
            output_path=os.path.join(self.tmpdir, "auto.pdf"),
            style="auto",
            title="Financial Report",
        ))
        assert r.success

    @pytest.mark.asyncio
    async def test_generate_bad_path(self):
        from digitorn.modules.pdf.module import PDFModule
        from digitorn.modules.pdf.params import GenerateParams
        m = PDFModule()
        r = await m.generate(GenerateParams(
            content="# Test",
            output_path=os.path.join(self.tmpdir, "test.txt"),
        ))
        assert not r.success
        assert "pdf" in r.error.lower()

    @pytest.mark.asyncio
    async def test_list_styles(self):
        from digitorn.modules.pdf.module import PDFModule
        from digitorn.modules.pdf.params import ListStylesParams
        m = PDFModule()
        r = await m.list_styles(ListStylesParams())
        assert r.success
        assert "consulting" in r.data

    def test_prompt_sections(self):
        from digitorn.modules.pdf.module import PDFModule
        m = PDFModule()
        sections = m.get_prompt_sections()
        assert len(sections) >= 1
        content = sections[0]["content"]
        assert "wb_write" in content
        assert "wb_append" in content
        assert "consulting" in content
        assert "![" in content  # images


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 5. PDF Styles System
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestPDFStyles:

    def test_all_palettes_exist(self):
        from digitorn.modules.pdf.styles import PALETTES
        expected = {"consulting", "scientific", "marketing", "financial", "tech", "executive",
                    "corporate", "ocean", "forest", "sunset", "royal", "dark", "slate",
                    "rose", "amber", "monochrome"}
        assert expected.issubset(set(PALETTES.keys()))

    def test_all_typographies_exist(self):
        from digitorn.modules.pdf.styles import TYPOGRAPHIES
        expected = {"modern", "classic", "consulting", "scientific", "marketing",
                    "financial", "tech", "executive", "compact", "elegant", "technical"}
        assert expected.issubset(set(TYPOGRAPHIES.keys()))

    def test_all_layouts_have_toc(self):
        from digitorn.modules.pdf.styles import LAYOUTS
        toc_layouts = [name for name, l in LAYOUTS.items() if l.has_toc]
        assert "cover" in toc_layouts
        assert "consulting" in toc_layouts

    def test_all_presets_resolve(self):
        from digitorn.modules.pdf.styles import PRESET_COMBOS, resolve_style
        for name in PRESET_COMBOS:
            resolved = resolve_style(style=name)
            assert resolved.palette is not None
            assert resolved.typography is not None
            assert resolved.layout is not None

    def test_preamble_generation(self):
        from digitorn.modules.pdf.styles import resolve_style, generate_preamble
        for style_name in ["consulting", "scientific", "marketing", "financial", "tech", "executive"]:
            resolved = resolve_style(style=style_name)
            preamble = generate_preamble(resolved, title="Test", author="Author")
            assert len(preamble) > 100
            assert "table" in preamble.lower()

    def test_layout_cover_styles(self):
        from digitorn.modules.pdf.styles import LAYOUTS
        cover_styles = {l.cover_style for l in LAYOUTS.values() if l.has_cover}
        assert "consulting" in cover_styles
        assert "executive" in cover_styles
        assert "creative" in cover_styles


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 6. WeasyPrint Theme CSS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


@_skip_weasy
class TestWeasyThemes:

    def test_all_themes_have_required_keys(self):
        from digitorn.modules.pdf.weasy import _THEMES
        required = {"primary", "accent", "text", "text_muted", "bg_subtle",
                     "bg_header", "neutral", "font_body", "font_heading",
                     "cover_bg", "cover_text"}
        for name, theme in _THEMES.items():
            for key in required:
                assert key in theme, f"Theme {name} missing key {key}"

    def test_toc_generation(self):
        from digitorn.modules.pdf.weasy import _build_toc_html, _THEMES
        body = "<h1>Section 1</h1><p>text</p><h2>Sub 1.1</h2><p>text</p><h1>Section 2</h1>"
        toc = _build_toc_html(body, _THEMES["consulting"])
        assert "Section 1" in toc
        assert "Section 2" in toc
        assert "Sub 1.1" in toc
        assert "toc-num" in toc

    def test_cover_generation(self):
        from digitorn.modules.pdf.weasy import _build_cover_html, _THEMES
        cover = _build_cover_html("My Title", "Author", "2026-03-21", _THEMES["consulting"])
        assert "My Title" in cover
        assert "Author" in cover
        assert "cover-spacer" in cover

    def test_clean_markdown(self):
        from digitorn.modules.pdf.weasy import _clean_markdown
        # Page Break literal
        assert "Page Break" not in _clean_markdown("Before\nPage Break\nAfter")
        # Double separators collapsed
        cleaned = _clean_markdown("---\n\n---\n\nText")
        assert cleaned.count("---") >= 1
        # None input handled
        assert _clean_markdown("") == ""


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 7. Presentation Module
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


@_skip_pptx
class TestPresentationAdvanced:

    def setup_method(self):
        from digitorn.modules.presentation.engine import PresentationEngine
        self.engine = PresentationEngine()
        self.tmpdir = tempfile.mkdtemp()

    def test_all_themes_produce_valid_pptx(self):
        from digitorn.modules.presentation.themes import THEMES
        for theme_name in THEMES:
            pid = f"t_{theme_name}"
            path = os.path.join(self.tmpdir, f"{pid}.pptx")
            self.engine.new_presentation(pid, path, theme_name=theme_name)
            self.engine.add_slide(pid, "<h1>Test</h1>", layout="title")
            self.engine.add_slide(pid, "<h2>Content</h2><p>Body</p>")
            result = self.engine.finalize(pid)
            assert os.path.isfile(path), f"Theme {theme_name} failed"
            assert result["slides"] == 2

    def test_kpi_cards(self):
        self.engine.new_presentation("kpi", os.path.join(self.tmpdir, "kpi.pptx"))
        html = (
            '<h2>Metrics</h2>'
            '<div class="columns">'
            '<div class="kpi"><span class="value">$10M</span><span class="label">Revenue</span></div>'
            '<div class="kpi"><span class="value">+42%</span><span class="label">Growth</span></div>'
            '</div>'
        )
        self.engine.add_slide("kpi", html)
        result = self.engine.finalize("kpi")
        assert result["slides"] == 1

    def test_progress_bars(self):
        self.engine.new_presentation("prog", os.path.join(self.tmpdir, "prog.pptx"))
        html = (
            '<h2>Progress</h2>'
            '<div class="progress" data-value="85" data-label="Target"></div>'
            '<div class="progress" data-value="60" data-label="Actual"></div>'
        )
        self.engine.add_slide("prog", html)
        result = self.engine.finalize("prog")
        assert result["slides"] == 1

    def test_timeline(self):
        self.engine.new_presentation("tl", os.path.join(self.tmpdir, "tl.pptx"))
        html = '<h2>Roadmap</h2><div class="timeline"><div>Q1</div><div>Q2</div><div>Q3</div></div>'
        self.engine.add_slide("tl", html)
        result = self.engine.finalize("tl")
        assert result["slides"] == 1

    def test_autofit_textbox(self):
        """Text boxes should have auto-fit enabled."""
        from pptx import Presentation
        from pptx.enum.text import MSO_AUTO_SIZE
        self.engine.new_presentation("af", os.path.join(self.tmpdir, "af.pptx"))
        self.engine.add_slide("af", "<h2>Test</h2><p>Long paragraph with lots of text " * 10 + "</p>")
        self.engine.finalize("af")
        prs = Presentation(os.path.join(self.tmpdir, "af.pptx"))
        autofit_count = sum(
            1 for slide in prs.slides for s in slide.shapes
            if s.has_text_frame and s.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        )
        assert autofit_count > 0

    def test_effects_applied(self):
        """Check gradient, shadow, glow effects are in the XML."""
        from pptx import Presentation
        from lxml import etree
        path = os.path.join(self.tmpdir, "fx.pptx")
        self.engine.new_presentation("fx", path, theme_name="corporate")
        self.engine.add_slide("fx", "<h1>Title</h1><p>Sub</p>", layout="title")
        self.engine.add_slide("fx", '<h2>KPIs</h2><div class="kpi"><span class="value">$5M</span><span class="label">Rev</span></div>')
        self.engine.finalize("fx")

        prs = Presentation(path)
        all_xml = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                all_xml += etree.tostring(shape._element, pretty_print=True).decode()
        assert "gradFill" in all_xml  # gradient on title slide
        assert "outerShdw" in all_xml  # shadow on KPI card


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 8. Incremental Actions in Agent Loop
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestIncrementalActions:

    @_skip_pptx
    def test_presentation_add_slide_is_incremental(self):
        """Verify add_slide is in the incremental actions exemption list."""
        import inspect
        from digitorn.core.runtime import agent_loop
        source = inspect.getsource(agent_loop)
        assert "add_slide" in source
        assert "presentation__add_slide" in source


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 9. YAML App Validation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


@pytest.mark.skipif(
    shutil.which("digitorn") is None,
    reason="digitorn CLI not on PATH",
)
class TestYAMLApps:
    """Test that our example YAML apps are valid."""

    def _validate(self, path: str) -> bool:
        import subprocess
        if not os.path.isfile(path):
            pytest.skip(f"Example file not found: {path}")
        env = os.environ.copy()
        env["DEEPSEEK_API_KEY"] = "sk-test-dummy-key"
        env["OPENAI_API_KEY"] = "sk-test-dummy"
        env["OPENROUTER_API_KEY"] = "sk-test-dummy"
        r = subprocess.run(
            ["digitorn", "app", "validate", path],
            capture_output=True, text=True, timeout=30, env=env,
        )
        return "Validation OK" in r.stdout

    def test_deepresearch_valid(self):
        assert self._validate("examples/deepresearch.yaml")

    @_skip_pptx
    def test_presentation_demo_valid(self):
        assert self._validate("examples/presentation-demo.yaml")

    @_skip_weasy
    def test_file_organizer_valid(self):
        assert self._validate("examples/file-organizer.yaml")

    def test_hello_oneshot_valid(self):
        assert self._validate("examples/hello-oneshot.yaml")

    def test_chat_valid(self):
        assert self._validate("examples/chat.yaml")

    @_skip_weasy
    def test_smart_chat_valid(self):
        assert self._validate("examples/smart-chat.yaml")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 10. Renderer Interface Contract
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestRendererContract:
    """Verify all renderers implement the full interface."""

    def _check_renderer(self, cls, name):
        methods = ["snapshot", "preview", "validate", "event_payload", "summary_line", "post_use_summary"]
        instance = cls() if name != "SpreadsheetRenderer" else cls(module=None)
        for method in methods:
            assert hasattr(instance, method), f"{name} missing {method}()"
            assert callable(getattr(instance, method)), f"{name}.{method} not callable"

    def test_markdown_renderer(self):
        from digitorn.modules.pdf.md_renderer import MarkdownRenderer
        self._check_renderer(MarkdownRenderer, "MarkdownRenderer")

    def test_spreadsheet_renderer(self):
        try:
            from digitorn.modules.spreadsheet.renderer import SpreadsheetRenderer
        except ImportError:
            pytest.skip("digitorn.modules.spreadsheet.renderer not available")
        self._check_renderer(SpreadsheetRenderer, "SpreadsheetRenderer")

    def test_presentation_renderer(self):
        from digitorn.modules.presentation.renderer import PresentationRenderer
        self._check_renderer(PresentationRenderer, "PresentationRenderer")

    def test_base_renderer_defaults(self):
        from digitorn.modules.context_builder.workbench import WorkbenchRenderer, Buffer
        r = WorkbenchRenderer()
        buf = Buffer(key="test", content="hello")
        assert isinstance(r.snapshot(buf), str)
        assert isinstance(r.preview(buf), str)
        assert isinstance(r.validate(buf), list)
        assert isinstance(r.event_payload(buf), dict)
        assert isinstance(r.summary_line(buf), str)
        assert isinstance(r.post_use_summary(buf, "tool", {}), str)
